from matplotlib import pyplot as plt
import pandas as pd
from pptx import Presentation
import os
import PySimpleGUI as sg
from pptx.util import Pt


def menu():
    sg.theme('Dark Blue 3')
    layout = [
              [sg.Text('Caminho do arquivo de Gestão do backlog')],
              [sg.Input(), sg.FileBrowse(key='plan', file_types=(('Text Files', '*.xls'),
                                                                 ('Text Files', '*.xlsx')))],
              [sg.Text('Nome do Projeto (conforme planilha)')],
              [sg.InputText(key='projeto')],
              [sg.Text('Sprint que deseja fechar')],
              [sg.InputText(key='sprint')],
              [sg.Text('Horas contratado por Sprint')],
              [sg.InputText(key='horas_sprint')],
              [[sg.Text('Caminho da pasta onde quer salvar os gráficos, \nrelatório e onde está o logo do cliente')],
               [sg.Input(), sg.FolderBrowse(key='path_save')],
               [[sg.Text('Caminho do Template PPT')],
               [sg.Input(), sg.FileBrowse(key='path_template')],
               [sg.Button('Gerar relatório'), sg.Button('Cancelar')]]]]
    return sg.Window('Geração de RFS', layout=layout, finalize=True,)


def erro():
    sg.theme('DarkRed')
    layout = [[sg.Text('Favor verificar:\n'
                       '\n1) Extensão do arquivo template em .pptx;\n'
                       '2) Extensão do arquivo de gestão em .xlsx;\n'
                       '3) Logo está na pasta selecionada;\n'
                       '4) Existe a Sprint na planilha;\n'
                       '5) Existe o cliente na planilha.')],
              [sg.Button('Voltar'), sg.Button('Cancelar')]]
    return sg.Window('ERRO', layout=layout, size=(500, 200), finalize=True)


def sucesso():
    sg.theme('DarkGreen')
    layout = [[sg.Text('Relatório gerado com sucesso !')],
              [sg.Button('Voltar'), sg.Button('Cancelar')]]
    return sg.Window('SUCESSO', layout=layout, size=(300, 100), finalize=True)


janela1, janela2, janela3 = menu(), None, None
while True:
    window, event, values = sg.read_all_windows()

# Operações no MENU
    if window == janela1 and event == sg.WINDOW_CLOSED:
        break
    if window == janela1 and event == 'Cancelar':
        break
    if window == janela1 and event == 'Gerar relatório' \
            and (values['projeto'] != '' and values['sprint'] != ''
                 and values['path_save'] != '' and values['plan'] != '' and
                 values['path_template'] != '' and values['horas_sprint'] != ''):
        try:
            plan = values['plan']
            sprint = values['sprint']
            projeto = values['projeto']
            save = values['path_save']
            template = values['path_template']
            plan_sprint = int(sprint) + 1
            horas_sprint = int(values['horas_sprint'])
            df = pd.DataFrame(pd.read_excel(plan, sheet_name='Sheets'))

# BurnUP
            total_sprint = int(sprint) + 1
            x = range(total_sprint)
            y = []

    # Horas contrato
            indice = 1
            hora_contrato = 0
            while indice < total_sprint:
                hora_contrato = hora_contrato + horas_sprint
                y.append(hora_contrato)
                indice = indice + 1
            y.append(0)
            y.sort()

    # Horas acumulado
            list_teste = []
            horas_acumulado = 0
            teste = []
            teste.append(0)
            serie = df[(df['Projeto'] == projeto) & (df['Sprint'])]
            z = list(round(serie.groupby(['Sprint'])['Tempo Efetivo'].sum(), 2))
            for i in z:
                horas_acumulado = horas_acumulado + i
                teste.append(horas_acumulado)
                if len(teste) > int(sprint):
                    break
            z = teste
            plt.plot(x, y, z)
            for x, t in zip(x, z):
                label = round(t, 2)
                plt.annotate(label, (x, t), textcoords="offset points", xytext=(0, 5), ha='center')
            plt.legend(['Contratado', 'Executado'])
            plt.title('Horas Executadas por Sprint', loc='center')
            plt.ylabel('HORAS EXECUTADAS')
            plt.xlabel('SPRINT')
            fig = plt.gcf()
            fig.set_size_inches((13.66, 7.68), forward=False)
            plt.savefig(save + '/BurnUp' + '.png', bbox_inches='tight', dpi=100)
            plt.show()

# FECHAMENTO DA SPRINT

    # Saldo
            saldo = float(hora_contrato) - float(horas_acumulado)
            if saldo < 0:
                texto_cons = 'Ao fim da Sprint ' + str(sprint) + ' a ' + projeto + ' tem um saldo negativo de ' + str(
                    (round(saldo), 2)) + ' Horas com a Indicium.'
                texto_cons2 = 'É natural que a Indicium realize mais horas conforme o avanço fluido do projeto.'
            else:
                texto_cons = 'Ao fim da Sprint ' + str(sprint) + ' a ' + projeto + ' tem um saldo positivo de ' + str(
                    (round(saldo), 2)) + ' Horas com a Indicium.'
                texto_cons2 = 'O saldo será consumido ao longo das demais Sprints.'

    # Fechamento maior consumo de horas na entrega
            serie = df[(df['Sprint'] == str(sprint).zfill(2)) & (df['Projeto'] == projeto)]
            total_horas = round(serie.groupby(['Sprint'])['Tempo Efetivo'].sum().max(), 2)
            agrupamento = serie.groupby(['Frente de Trabalho'])['Tempo Efetivo'].sum()
            maior_entrega = agrupamento.idxmax()
            horas_entrega = round(agrupamento.max(), 2)
            texto_cons_entrega = 'A maior parte das horas consumidas na Sprint foram destinadas a: ' + maior_entrega + \
                                 '. Ela representa um total de ' + str(horas_entrega) + ' horas consumidas de um total de ' \
                                 + str(total_horas) + ' horas executadas'

    # Fechamento maior consumo de horas pessoa
            serie = df[(df['Sprint'] == str(sprint).zfill(2)) & (df['Projeto'] == projeto)]
            agrupamento = serie.groupby(['Responsavel'])['Tempo Efetivo'].sum()
            maior_responsavel = agrupamento.idxmax()
            horas_responsavel = round(agrupamento.max(), 2)
            texto_cons_resp = 'A maior parte das horas consumidas na Sprint foram de responsabilidade do(a): ' + \
                              maior_responsavel + '. Ela representa um total de ' + str(horas_responsavel) + \
                              ' horas consumidas de um total de ' + str(total_horas) + ' horas executadas'

    # Planejamento maior consumo de horas na entrega
            serie = df[(df['Sprint'] == str(plan_sprint).zfill(2)) & (df['Projeto'] == projeto)]
            total_horas = round(serie.groupby(['Sprint'])['Tempo Estimado'].sum().max(), 2)
            agrupamento = serie.groupby(['Frente de Trabalho'])['Tempo Estimado'].sum()
            maior_entrega = agrupamento.idxmax()
            horas_entrega = round(agrupamento.max(), 2)
            texto_plan_ent = 'A maior parte das horas planejadas na Sprint serão destinadas a: ' + maior_entrega + \
                             '. Ela representa um total de ' + str(horas_entrega) + ' horas consumidas de um total de ' \
                             + str(total_horas) + ' horas planejadas'

    # Planejamento maior consumo de horas pessoa
            serie = df[(df['Sprint'] == str(plan_sprint).zfill(2)) & (df['Projeto'] == projeto)]
            agrupamento = serie.groupby(['Responsavel'])['Tempo Estimado'].sum()
            maior_responsavel = agrupamento.idxmax()
            horas_responsavel = round(agrupamento.max(), 2)
            texto_plan_resp = 'A maior parte das horas planejadas na Sprint serão de responsabilidade do(a): ' + maior_responsavel \
                              + '. Ela representa um total de ' + str(horas_responsavel) + \
                              ' horas consumidas de um total de ' + str(total_horas) + ' horas planejadas'

    # Fechamento - Gráfico de horas por Entrega
            serie = df[(df['Sprint'] == str(sprint).zfill(2)) & (df['Projeto'] == projeto)]
            agrupamento = serie.groupby(['Frente de Trabalho'])['Tempo Efetivo'].sum().reset_index()
            x = agrupamento['Frente de Trabalho']
            y = agrupamento['Tempo Efetivo']
            plt.bar(x, y)
            for x, y in zip(x, y):
                label = round(y, 2)
                plt.annotate(label, (x, y), textcoords="offset points", xytext=(0, 5), ha='center')
            plt.title('Horas Executadas por Entrega - Sprint ' + str(sprint), loc='center')
            plt.ylabel('HORAS EXECUTADAS')
            plt.xlabel('ENTREGAS')
            fig = plt.gcf()
            fig.set_size_inches((13.66, 7.68), forward=False)
            plt.savefig(save + '/Horas Executadas por Entrega - Sprint ' + str(sprint) + '.png',
                        bbox_inches='tight',
                        dpi=100)
            plt.show()

    # Fechamento - Gráfico de horas por Responsável
            serie = df[(df['Sprint'] == str(sprint).zfill(2)) & (df['Projeto'] == projeto)]
            agrupamento = serie.groupby(['Responsavel'])['Tempo Efetivo'].sum().reset_index()
            x = agrupamento['Responsavel']
            y = agrupamento['Tempo Efetivo']
            plt.bar(x, y)
            for x, y in zip(x, y):
                label = round(y, 2)
                plt.annotate(label, (x, y), textcoords="offset points", xytext=(0, 5), ha='center')
            plt.title('Horas Executadas por Responsável - Sprint ' + str(sprint), loc='center')
            plt.ylabel('HORAS EXECUTADAS')
            plt.xlabel('RESPONSÁVEL')
            fig = plt.gcf()
            fig.set_size_inches((13.66, 7.68), forward=False)
            plt.savefig(save + '/Horas Executadas por Responsável - Sprint ' + str(sprint) + '.png',
                        bbox_inches='tight', dpi=100)
            plt.show()

# PLANEJAMENTO DA SPRINT

    # Planejamento - Gráfico de horas por Responsável
            serie = df[(df['Sprint'] == str(plan_sprint).zfill(2)) & (df['Projeto'] == projeto)]
            agrupamento = serie.groupby(['Responsavel'])['Tempo Estimado'].sum().reset_index()
            x = agrupamento['Responsavel']
            y = agrupamento['Tempo Estimado']
            plt.bar(x, y)
            for x, y in zip(x, y):
                label = round(y, 2)
                plt.annotate(label, (x, y), textcoords="offset points", xytext=(0, 5), ha='center')
            plt.title('Horas Planejadas por Responsável - Sprint ' + str(plan_sprint), loc='center')
            plt.ylabel('HORAS PLANEJADAS')
            plt.xlabel('RESPONSÁVEL')
            fig = plt.gcf()
            fig.set_size_inches((13.66, 7.68), forward=False)
            plt.savefig(save + '/Horas Planejadas por Responsável - Sprint ' + str(plan_sprint) + '.png',
                        bbox_inches='tight', dpi=100)
            plt.show()

    # Planejamento - Gráfico de horas por Entrega
            serie = df[(df['Sprint'] == str(plan_sprint).zfill(2)) & (df['Projeto'] == projeto)]
            agrupamento = serie.groupby(['Frente de Trabalho'])['Tempo Estimado'].sum().reset_index()
            x = agrupamento['Frente de Trabalho']
            y = agrupamento['Tempo Estimado']
            plt.bar(x, y)
            for x, y in zip(x, y):
                label = round(y, 2)
                plt.annotate(label, (x, y), textcoords="offset points", xytext=(0, 5), ha='center')
            plt.title('Horas Planejadas por Entrega - Sprint ' + str(plan_sprint), loc='center')
            plt.ylabel('HORAS PLANEJADAS')
            plt.xlabel('ENTREGAS')
            fig = plt.gcf()
            fig.set_size_inches((13.66, 7.68), forward=False)
            plt.savefig(save + '/Horas Planejadas por Entrega - Sprint ' + str(plan_sprint) + '.png',
                        bbox_inches='tight', dpi=100)
            plt.show()

# Salvar grafico no ppt
            prs = Presentation(template)
            logo = save + '/' + projeto + '.png'

    # Logo do cliente
            slide = prs.slides[1]
            picture_logo = slide.shapes.add_picture(logo, left=800000, top=2000000, width=None, height=None)

    # BurnUp
            img_burnup = save + '/BurnUp' + '.png'
            slide = prs.slides[4]
            picture_burnup = slide.shapes.add_picture(img_burnup, left=300000, top=1100000, width=5000000, height=2800000)
            text_ph1 = slide.placeholders[1]
            text_ph1.text_frame.text = texto_cons + ' ' + texto_cons2
            font = text_ph1.text_frame.paragraphs[0].runs[0].font
            font.name = 'Roboto'
            font.size = Pt(11)
            font.bold = False
            font.italic = False

    # Fechamento - Entregas
            img_entregas = save + '/Horas Executadas por Entrega - Sprint ' + str(sprint) + '.png'
            slide = prs.slides[5]
            picture_entregas = slide.shapes.add_picture(img_entregas, left=500000, top=1300000, width=5000000,
                                                        height=3100000)
            text_ph1 = slide.placeholders[1]
            text_ph1.text_frame.text = texto_cons_entrega
            font = text_ph1.text_frame.paragraphs[0].runs[0].font
            font.name = 'Roboto'
            font.size = Pt(11)
            font.bold = False
            font.italic = False

    # Fechamento - Responsável
            img_responsavel = save + '/Horas Executadas por Responsável - Sprint ' + str(sprint) + '.png'
            slide = prs.slides[6]
            picture_responsavel = slide.shapes.add_picture(img_responsavel, left=500000, top=1300000, width=5000000,
                                                           height=3100000)
            text_ph1 = slide.placeholders[1]
            text_ph1.text_frame.text = texto_cons_resp
            font = text_ph1.text_frame.paragraphs[0].runs[0].font
            font.name = 'Roboto'
            font.size = Pt(11)
            font.bold = False
            font.italic = False

    # Planejamento - Entregas
            img_plan_entrega = save + '/Horas Planejadas por Entrega - Sprint ' + str(plan_sprint) + '.png'
            slide = prs.slides[7]
            picture_plan_entrega = slide.shapes.add_picture(img_plan_entrega, left=500000, top=1300000, width=5000000,
                                                            height=3100000)
            text_ph1 = slide.placeholders[1]
            text_ph1.text_frame.text = texto_plan_ent
            font = text_ph1.text_frame.paragraphs[0].runs[0].font
            font.name = 'Roboto'
            font.size = Pt(11)
            font.bold = False
            font.italic = False

    # Planejamento - Responsável
            img_plan_responsavel = save + '/Horas Planejadas por Responsável - Sprint ' + str(
                plan_sprint) + '.png'
            slide = prs.slides[8]
            picture_plan_responsavel = slide.shapes.add_picture(img_plan_responsavel, left=500000, top=1300000,
                                                                width=5000000, height=3100000)
            text_ph1 = slide.placeholders[1]
            text_ph1.text_frame.text = texto_plan_resp
            font = text_ph1.text_frame.paragraphs[0].runs[0].font
            font.name = 'Roboto'
            font.size = Pt(11)
            font.bold = False
            font.italic = False

# Salvar arquivo
            prs.save(save + "/" + projeto + ' - Relatório de Fechamento de Sprint ' + sprint + ".pptx")
            os.startfile(save + "/" + projeto + ' - Relatório de Fechamento de Sprint ' + sprint + ".pptx")
            janela1.close()
            janela3 = sucesso()
        except:
            janela1.close()
            janela2 = erro()

# Operações Janela ERRO
    if window == janela2 and event == 'Voltar':
        janela2.close()
        janela1 = menu()
    if window == janela2 and event == 'Cancelar':
        break
    if window == janela2 and event == sg.WINDOW_CLOSED:
        break

# Operações Janela SUCESSO
    if window == janela3 and event == 'Voltar':
        janela3.close()
        janela1 = menu()
    if window == janela3 and event == 'Cancelar':
        break
    if window == janela3 and event == sg.WINDOW_CLOSED:
        break
